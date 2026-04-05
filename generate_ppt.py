from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_ppt():
    prs = Presentation()
    
    # helper for slide creation
    def add_slide(title_text, content_text, sub_headings=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Access elements
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        # Set White Background (Default is white, but let's be explicit if needed)
        # slide.background.fill.solid()
        # slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Title Formatting
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.alignment = PP_ALIGN.LEFT
            
        # Body Formatting
        tf = body.text_frame
        tf.clear() # clear default bullet
        
        if isinstance(content_text, list):
            for i, p_info in enumerate(content_text):
                p = tf.add_paragraph()
                if isinstance(p_info, dict):
                    # For subheadings or specialized matter
                    p.text = p_info.get('text', '')
                    p.font.size = Pt(p_info.get('size', 12))
                    p.font.bold = p_info.get('bold', False)
                else:
                    p.text = str(p_info)
                    p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_after = Pt(10)
        else:
            p = tf.add_paragraph()
            p.text = content_text
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0, 0, 0)

    # 1. Slide Title
    add_slide("Yukti: ImpactSim", [
        {"text": "AI-Powered Social Entrepreneurship Simulator", "size": 14, "bold": True},
        {"text": "Empowering the next generation of social innovators through High-Performance AI.", "size": 12}
    ])

    # 2. Abstract
    add_slide("Abstract", [
        "ImpactSim (Yukti) is an immersive, AI-first simulation platform designed to bridge the gap between theoretical social entrepreneurship and real-world execution. The system provides a 'safe space' for youth to practice decision-making, resource management, and stakeholder navigation. By utilizing Llama 3.1 via Groq's high-speed inference, the platform offers real-time mentoring and dynamic boardroom simulations, enabling users to transform vague social concepts into actionable, validated enterprise models."
    ])

    # 3. Introduction
    add_slide("Introduction", [
        {"text": "Overview", "size": 14, "bold": True},
        "Social entrepreneurship is critical for solving global challenges, yet many aspiring changemakers are deterred by the high risk of failure and lack of structured mentorship. Yukti provides a structured pathway for these individuals.",
        {"text": "Key Focus Areas", "size": 14, "bold": True},
        "- Gamified Idea Development (Sprint Game)",
        "- Interactive AI-Driven Boardroom Simulations powered by Llama 3.1",
        "- Immersive 3D/Neural User Interfaces"
    ])

    # 4. Existing System
    add_slide("Existing System", [
        {"text": "Current Limitations", "size": 14, "bold": True},
        "- Reliance on static case studies and theoretical textbooks.",
        "- High cost and low accessibility of physical incubation programs.",
        "- Lack of real-time, personalized feedback on startup ideas.",
        "- Non-interactive learning models that fail to simulate market dynamics."
    ])

    # 5. Proposed System
    add_slide("Proposed System", [
        {"text": "Innovations in Yukti", "size": 14, "bold": True},
        "- Visionary Architect Simulation Engine (VASE): A custom Groq-powered AI engine for multi-step validation.",
        "- AI Boardroom: Real-time stakeholder negotiation simulator powered by Llama 3.1.",
        "- Neural Hub: Centralized dashboard for tracking environmental and social impact metrics.",
        "- Multi-language Mentorship: Breaking barriers with integrated Hindi and English support."
    ])

    # 6. Workflow
    add_slide("Workflow", [
        {"text": "Process Cycle", "size": 14, "bold": True},
        "1. Idea Intake: User captures a social problem and propose a solution.",
        "2. The Sprint Game: 5-step gamified journey to refine the business model.",
        "3. AI Validation: Llama 3.1-driven analysis of feasibility and impact.",
        "4. Boardroom Simulation: Presentation to virtual stakeholders and handling rebuttals.",
        "5. Iteration: Updating the Workspace based on simulation feedback."
    ])

    # 7. System Architecture
    add_slide("System Architecture", [
        {"text": "Core Components", "size": 14, "bold": True},
        "- Frontend: React 18, Vite, Tailwind CSS, Three.js (for Neural Flux Field).",
        "- Backend & Auth: Firebase Authentication and Firestore Database.",
        "- AI Integration: Groq Cloud SDK with Llama 3.1-8b-instant models.",
        "- Routing & State: React Context API and Animate Presence for transitions."
    ])

    # 8. System Requirements
    add_slide("System Requirements", [
        {"text": "Hardware", "size": 14, "bold": True},
        "- Any device with a modern web browser (Chrome, Edge, Safari).",
        "- Minimum 4GB RAM for smooth 3D rendering.",
        {"text": "Software", "size": 14, "bold": True},
        "- OS: Windows, Linux, or macOS.",
        "- Environment: Node.js (v18+) for local execution.",
        "- Stable Internet Connection for AI and Firebase services."
    ])

    # 9. Reference
    add_slide("Reference", [
        "- Groq Cloud API Documentation (console.groq.com)",
        "- Meta Llama 3.1 Model Card (ai.meta.com/llama)",
        "- Firebase Documentation (firebase.google.com/docs)",
        "- React.js Official Documentation (react.dev)"
    ])

    prs.save('Yukti_ImpactSim_Presentation.pptx')
    print("PPT generated successfully!")

if __name__ == "__main__":
    create_ppt()
